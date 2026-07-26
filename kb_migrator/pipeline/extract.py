"""文本提取 / 格式解析。

按后缀分派解析器；解析库缺失或旧格式(.doc/.ppt/.xls)时优雅降级：
- 旧二进制格式建议先用 LibreOffice headless 转 docx/xlsx（此处给出转换钩子）；
- 纯文本按 charset-normalizer 检测编码，避免 GBK/GB2312 乱码；
- 提取失败不抛，返回空字符串 + 原因，交由上层记 error_detail。

所有解析库均为可选依赖：缺失时该类型跳过而非崩溃，保证 MVP 可跑。
"""
from __future__ import annotations

import os
import subprocess
import sys
import tempfile
from dataclasses import dataclass
from pathlib import Path


@dataclass
class ExtractResult:
    text: str
    ok: bool
    note: str = ""


def _read_text_auto(path: str) -> str:
    """用 charset-normalizer 检测编码读纯文本；失败退回 utf-8/ignore。"""
    try:
        from charset_normalizer import from_path

        best = from_path(path).best()
        if best is not None:
            return str(best)
    except Exception:
        pass
    with open(path, "rb") as f:
        return f.read().decode("utf-8", errors="ignore")


def _extract_docx(path: str) -> str:
    from docx import Document  # python-docx

    doc = Document(path)
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            parts.append("\t".join(c.text for c in row.cells))
    return "\n".join(parts)


def _extract_xlsx(path: str) -> str:
    from openpyxl import load_workbook

    wb = load_workbook(path, read_only=True, data_only=True)
    lines: list[str] = []
    for ws in wb.worksheets:
        lines.append(f"# sheet: {ws.title}")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                lines.append("\t".join(cells))
    wb.close()
    return "\n".join(lines)


def _extract_pptx(path: str) -> str:
    from pptx import Presentation

    prs = Presentation(path)
    parts: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        parts.append(f"# slide {i}")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = "".join(r.text for r in para.runs)
                    if t.strip():
                        parts.append(t)
    return "\n".join(parts)


def _extract_pdf(path: str) -> str:
    # 优先 pdfplumber（版面/表格更好），退回 pypdf
    try:
        import pdfplumber

        out = []
        with pdfplumber.open(path) as pdf:
            for page in pdf.pages:
                out.append(page.extract_text() or "")
        text = "\n".join(out).strip()
        if text:
            return text
    except Exception:
        pass
    try:
        from pypdf import PdfReader

        reader = PdfReader(path)
        return "\n".join((p.extract_text() or "") for p in reader.pages)
    except Exception:
        return ""


def _libreoffice_convert(path: str, target_ext: str) -> str | None:
    """用 LibreOffice headless 把旧格式转成现代格式，返回新文件路径。

    需系统装有 soffice；未装则返回 None（上层降级为跳过并告警）。
    跨平台兜底路径：Windows 上若装了 Office，优先走 _com_convert。
    """
    soffice = os.environ.get("KBM_SOFFICE", "soffice")
    outdir = tempfile.mkdtemp(prefix="kbm_lo_")
    try:
        subprocess.run(
            [soffice, "--headless", "--convert-to", target_ext, "--outdir", outdir, path],
            check=True, capture_output=True, timeout=120,
        )
    except (FileNotFoundError, subprocess.SubprocessError):
        return None
    stem = Path(path).stem
    candidate = os.path.join(outdir, f"{stem}.{target_ext}")
    return candidate if os.path.exists(candidate) else None


# Office COM SaveAs 文件格式常量
_WD_FORMAT_DOCX = 16   # wdFormatXMLDocument
_XL_FORMAT_XLSX = 51   # xlOpenXMLWorkbook
_PP_FORMAT_PPTX = 24   # ppSaveAsOpenXMLPresentation


def _com_convert(path: str, target_ext: str) -> str | None:
    """Windows 上用已安装的 Office(Word/Excel/PowerPoint) COM 把旧格式另存为现代格式。

    优点：企业机免安装、转换保真度高。无 pywin32 / 无 Office / 转换失败均返回 None，
    交由上层继续降级到 LibreOffice 或告警。始终确保退出应用，避免残留进程。
    """
    if sys.platform != "win32":
        return None
    try:
        import pythoncom  # noqa: F401
        import win32com.client
    except ImportError:
        return None

    src = os.path.abspath(path)
    outdir = tempfile.mkdtemp(prefix="kbm_com_")
    dst = os.path.join(outdir, f"{Path(path).stem}.{target_ext}")
    ext = os.path.splitext(path)[1].lower()

    # 每个线程调用 COM 需先 CoInitialize（Web 控制台在工作线程里跑）
    pythoncom.CoInitialize()
    app = None
    try:
        if ext == ".doc":
            app = win32com.client.DispatchEx("Word.Application")
            app.Visible = False
            app.DisplayAlerts = 0
            doc = app.Documents.Open(src, ConfirmConversions=False, ReadOnly=True)
            doc.SaveAs2(dst, FileFormat=_WD_FORMAT_DOCX)
            doc.Close(SaveChanges=False)
        elif ext == ".xls":
            app = win32com.client.DispatchEx("Excel.Application")
            app.Visible = False
            app.DisplayAlerts = False
            wb = app.Workbooks.Open(src, ReadOnly=True)
            wb.SaveAs(dst, FileFormat=_XL_FORMAT_XLSX)
            wb.Close(SaveChanges=False)
        elif ext == ".ppt":
            app = win32com.client.DispatchEx("PowerPoint.Application")
            # PowerPoint 不支持完全隐藏窗口，用 WithWindow=False 打开
            prs = app.Presentations.Open(src, ReadOnly=True, WithWindow=False)
            prs.SaveAs(dst, _PP_FORMAT_PPTX)
            prs.Close()
        else:
            return None
    except Exception:  # noqa: BLE001 —— COM 异常五花八门，统一降级
        return None
    finally:
        if app is not None:
            try:
                app.Quit()
            except Exception:  # noqa: BLE001
                pass
        pythoncom.CoUninitialize()
    return dst if os.path.exists(dst) else None


def _convert_legacy(path: str, target_ext: str) -> str | None:
    """旧格式转换：Windows 优先 Office COM（免安装、保真），否则 LibreOffice。"""
    return _com_convert(path, target_ext) or _libreoffice_convert(path, target_ext)


def _ocr_pdf(path: str, lang: str = "chi_sim+eng") -> tuple[str, str]:
    """对扫描件 PDF 逐页渲染成图再 Tesseract OCR。返回 (text, note)。

    依赖：PyMuPDF(渲染) + pytesseract + Pillow + 系统 Tesseract 可执行(含中文语言包)。
    Tesseract 路径可用 KBM_TESSERACT 指定；缺任一依赖/引擎则降级并给出原因。
    """
    try:
        import fitz  # PyMuPDF
        import pytesseract
        from PIL import Image
    except ImportError as e:
        return "", f"OCR 依赖缺失: {e.name}（需 pymupdf/pytesseract/Pillow）"

    cmd = os.environ.get("KBM_TESSERACT")
    if cmd:
        pytesseract.pytesseract.tesseract_cmd = cmd
    import io

    try:
        doc = fitz.open(path)
    except Exception as e:  # noqa: BLE001
        return "", f"OCR 打开 PDF 失败: {e}"
    parts: list[str] = []
    try:
        for page in doc:
            pix = page.get_pixmap(dpi=200)          # 200dpi 兼顾清晰度与速度
            img = Image.open(io.BytesIO(pix.tobytes("png")))
            parts.append(pytesseract.image_to_string(img, lang=lang))
    except pytesseract.TesseractNotFoundError:
        return "", "未找到 Tesseract 引擎（设 KBM_TESSERACT 指向 tesseract.exe 或加入 PATH）"
    except Exception as e:  # noqa: BLE001
        return "", f"OCR 失败: {e}"
    finally:
        doc.close()
    text = "\n".join(parts).strip()
    return (text, "") if text else ("", "OCR 未识别出文本")


_DISPATCH = {
    ".docx": _extract_docx,
    ".xlsx": _extract_xlsx,
    ".pptx": _extract_pptx,
    ".pdf": _extract_pdf,
}
_PLAIN = {".txt", ".md", ".csv", ".rtf"}
# 旧二进制格式 -> 转换目标
_LEGACY = {".doc": "docx", ".xls": "xlsx", ".ppt": "pptx"}


def extract_text(path: str) -> ExtractResult:
    """按后缀提取纯文本。缺依赖/旧格式/扫描件等场景优雅降级。"""
    ext = os.path.splitext(path)[1].lower()
    try:
        if ext in _PLAIN:
            return ExtractResult(_read_text_auto(path), ok=True)
        if ext in _DISPATCH:
            text = _DISPATCH[ext](path)
            if text.strip():
                return ExtractResult(text, ok=True)
            # PDF 无文本层 -> 可能是扫描件，尝试 OCR
            if ext == ".pdf":
                ocr_text, note = _ocr_pdf(path)
                if ocr_text.strip():
                    return ExtractResult(ocr_text, ok=True, note="OCR 提取")
                return ExtractResult("", ok=False, note=note or "empty_text (扫描件 OCR 无结果)")
            return ExtractResult("", ok=False, note="empty_text")
        if ext in _LEGACY:
            converted = _convert_legacy(path, _LEGACY[ext])
            if converted is None:
                return ExtractResult(
                    "", ok=False,
                    note="旧格式转换失败（需 Office COM 或 LibreOffice）")
            return extract_text(converted)
        return ExtractResult("", ok=False, note=f"unsupported ext {ext}")
    except ImportError as e:
        return ExtractResult("", ok=False, note=f"缺少解析依赖: {e.name}")
    except Exception as e:  # noqa: BLE001 —— 提取失败不应中断整批
        return ExtractResult("", ok=False, note=f"extract error: {e}")
